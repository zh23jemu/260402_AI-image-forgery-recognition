from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor


ROOT = Path(r"C:\Coding\260402_AI-image-forgery-recognition")
DOCX_PATH = ROOT / "AI伪造图像识别论文_终稿_改写2_改写.docx"
OUT_PPTX = ROOT / "docs" / "AI伪造图像识别论文_答辩PPT.pptx"
PREVIEW_DIR = ROOT / "analysis" / "defense_ppt_preview"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SLIDE_W = 1280
SLIDE_H = 720

NAVY = RGBColor(15, 35, 61)
TEAL = RGBColor(41, 173, 163)
BLUE = RGBColor(52, 121, 227)
ORANGE = RGBColor(244, 160, 72)
GREEN = RGBColor(51, 145, 89)
TEXT = RGBColor(26, 34, 48)
MUTED = RGBColor(95, 107, 123)
LIGHT_BG = RGBColor(248, 250, 252)
PALE = RGBColor(233, 239, 245)
LINE = RGBColor(214, 221, 231)
WHITE = RGBColor(255, 255, 255)
SOFT_BLUE = RGBColor(237, 244, 255)
SOFT_TEAL = RGBColor(235, 250, 248)
SOFT_ORANGE = RGBColor(255, 245, 232)

FONT_CJK = "Microsoft YaHei"
FONT_CJK_BOLD = "Microsoft YaHei"
FONT_LATIN = "Times New Roman"

FONT_FILE = Path(r"C:\Windows\Fonts\msyh.ttc")
FONT_FILE_BOLD = Path(r"C:\Windows\Fonts\simhei.ttf")
FONT_FILE_SERIF = Path(r"C:\Windows\Fonts\simsun.ttc")


@dataclass
class SlideSpec:
    title: str
    section: str


def in_to_px(x: float) -> int:
    return int(round(x * 96))


def px_to_in(px: int) -> float:
    return px / 96.0


def color_hex(rgb: RGBColor) -> str:
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def load_font(size: int, bold: bool = False, serif: bool = False) -> ImageFont.FreeTypeFont:
    candidates = []
    if serif:
        candidates.append(FONT_FILE_SERIF)
    if bold:
        candidates.append(FONT_FILE_BOLD)
    candidates.append(FONT_FILE)
    candidates.append(FONT_FILE_BOLD)
    for path in candidates:
        if path.exists():
            try:
                return ImageFont.truetype(str(path), size=size)
            except Exception:
                continue
    return ImageFont.load_default()


def set_slide_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 24,
    color: RGBColor = TEXT,
    bold: bool = False,
    font_face: str = FONT_CJK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    italic: bool = False,
    line_spacing: float | None = None,
    margins: tuple[float, float, float, float] = (0.04, 0.04, 0.02, 0.02),
):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margins[0])
    tf.margin_right = Inches(margins[1])
    tf.margin_top = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    tf.vertical_anchor = valign
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font_face
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tx


def add_multiline_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Sequence[str],
    *,
    font_size: int = 20,
    color: RGBColor = TEXT,
    bold: bool = False,
    font_face: str = FONT_CJK,
    bullet: bool = False,
):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        if bullet:
            p.bullet = True
        for run in p.runs:
            run.font.name = font_face
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tx


def add_rule(slide, left: float, top: float, width: float, color: RGBColor, weight: int = 2):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(weight / 96))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor | None = None, radius: bool = False, transparency: int = 0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    try:
        shp.fill.transparency = transparency
    except Exception:
        pass
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_chip(slide, left: float, top: float, text: str, fill: RGBColor, color: RGBColor = WHITE, width: float | None = None):
    width = width or max(0.75, 0.13 * len(text) + 0.4)
    shp = add_rect(slide, left, top, width, 0.34, fill, radius=True)
    tx = add_textbox(
        slide,
        left,
        top + 0.02,
        width,
        0.28,
        text,
        font_size=13,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margins=(0.02, 0.02, 0.0, 0.0),
    )
    return shp, tx


def add_footer(slide, page: int):
    add_textbox(
        slide,
        0.65,
        7.05,
        5.5,
        0.28,
        "AI伪造图像识别研究",
        font_size=10,
        color=MUTED,
        font_face=FONT_LATIN,
    )
    add_textbox(
        slide,
        12.1,
        7.02,
        0.45,
        0.3,
        str(page),
        font_size=12,
        color=MUTED,
        font_face=FONT_LATIN,
        align=PP_ALIGN.RIGHT,
    )


def add_content_header(slide, section: str, title: str, subtitle: str | None = None, page: int = 1):
    add_textbox(
        slide,
        0.65,
        0.25,
        4.3,
        0.25,
        section,
        font_size=12,
        color=TEAL,
        bold=True,
        font_face=FONT_CJK_BOLD,
    )
    add_textbox(
        slide,
        0.65,
        0.48,
        11.2,
        0.58,
        title,
        font_size=28,
        color=TEXT,
        bold=True,
        font_face=FONT_CJK_BOLD,
    )
    add_rule(slide, 0.65, 1.08, 1.15, TEAL, weight=4)
    if subtitle:
        add_textbox(
            slide,
            0.65,
            1.13,
            11.2,
            0.38,
            subtitle,
            font_size=14,
            color=MUTED,
            font_face=FONT_CJK,
        )
    add_footer(slide, page)


def add_panel(slide, left: float, top: float, width: float, height: float, fill: RGBColor = WHITE, line: RGBColor = LINE, radius: bool = True):
    return add_rect(slide, left, top, width, height, fill=fill, line=line, radius=radius)


def add_image(slide, path: Path, left: float, top: float, width: float, height: float):
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def build_pptx() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    # Slide 1 cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, NAVY)
    add_rect(slide, 0.0, 0.0, 13.333, 7.5, NAVY)
    add_rect(slide, 0.0, 0.0, 0.22, 7.5, ORANGE)
    add_rect(slide, 11.75, 0.0, 1.58, 7.5, RGBColor(11, 29, 52), transparency=0)
    add_rect(slide, 9.55, 0.85, 2.8, 5.85, RGBColor(19, 47, 82))
    add_rect(slide, 10.0, 1.2, 2.0, 0.2, TEAL)
    add_rect(slide, 10.0, 1.55, 1.6, 0.18, ORANGE)
    add_rect(slide, 10.0, 1.9, 1.2, 0.18, BLUE)
    add_textbox(slide, 0.72, 1.0, 8.2, 0.35, "答辩汇报", font_size=18, color=RGBColor(180, 205, 230), bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 0.72, 1.4, 8.7, 1.2, "AI伪造图像识别研究", font_size=36, color=WHITE, bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 0.72, 2.45, 7.9, 0.8, "面向泛化检测的多模型识别与联合验证", font_size=20, color=RGBColor(214, 231, 242), font_face=FONT_CJK)
    add_rule(slide, 0.72, 3.28, 2.1, TEAL, weight=4)
    add_textbox(slide, 0.72, 3.48, 6.9, 0.9, "从基线复现，到困难场景分析，再到两阶段联合验证", font_size=16, color=RGBColor(207, 219, 233), font_face=FONT_CJK)
    add_chip(slide, 0.72, 4.42, "FSD", TEAL, width=0.92)
    add_chip(slide, 1.78, 4.42, "Stay-Positive", ORANGE, width=1.7)
    add_chip(slide, 3.62, 4.42, "LVLM", BLUE, width=1.05)
    add_textbox(slide, 0.72, 5.4, 4.3, 1.15, "学生：待填写\n指导教师：待填写\n完成时间：2026年5月", font_size=14, color=RGBColor(215, 227, 239), font_face=FONT_CJK)
    add_textbox(slide, 10.05, 3.6, 1.95, 0.75, "多模型\n基线复现", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 10.05, 4.55, 1.95, 0.75, "困难场景\n样本分析", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 10.05, 5.5, 1.95, 0.75, "联合验证\n与总结", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font_face=FONT_CJK_BOLD)

    # Slide 2 background and problems
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT_BG)
    add_content_header(
        slide,
        "01 背景与问题",
        "为什么 AI 伪造图像检测越来越难",
        "生成器迭代快，单一检测器很难同时兼顾泛化性、稳定性与解释性。",
        page=2,
    )
    cols = [
        (0.72, TEAL, "01", "跨生成器泛化不足", [
            "训练时见过的生成器，和部署时遇到的生成器往往不一致。",
            "单一模型容易学到特定数据集或生成器的局部模式。",
        ]),
        (4.43, ORANGE, "02", "阈值与分数分布不稳定", [
            "真伪样本在边界区域容易重叠，0.5 附近最容易翻转。",
            "不同协议、不同生成器下的阈值会发生明显漂移。",
        ]),
        (8.14, BLUE, "03", "复杂样本难解释", [
            "误判样本往往不是“完全不像”，而是局部冲突和混合证据。",
            "只给出真假结果不够，还需要进一步说明为什么会错。",
        ]),
    ]
    for left, color, num, head, lines in cols:
        add_panel(slide, left, 1.65, 3.35, 4.4, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, left + 0.22, 1.85, 0.8, 0.45, num, font_size=26, color=color, bold=True, font_face=FONT_LATIN)
        add_textbox(slide, left + 0.22, 2.28, 2.8, 0.5, head, font_size=20, color=TEXT, bold=True, font_face=FONT_CJK_BOLD)
        add_rule(slide, left + 0.22, 2.86, 0.72, color, weight=3)
        add_multiline_text(slide, left + 0.22, 3.02, 2.85, 2.55, lines, font_size=16, color=TEXT, font_face=FONT_CJK)
    add_panel(slide, 0.72, 6.24, 11.9, 0.46, fill=SOFT_TEAL, line=SOFT_TEAL, radius=True)
    add_textbox(
        slide,
        1.0,
        6.32,
        11.3,
        0.3,
        "研究目标：构建一个面向泛化检测的多模型协同分析框架，把“复现、分析、联合验证”串成一条可复核的研究链条。",
        font_size=15,
        color=TEXT,
        bold=True,
        font_face=FONT_CJK,
    )

    # Slide 3 research flow and innovations
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_content_header(
        slide,
        "02 研究思路与创新点",
        "本文的研究主线与创新点",
        "不是只做一个结果，而是把复现、困难场景分析和联合验证连成闭环。",
        page=3,
    )
    # Flow strip
    x_positions = [0.78, 4.55, 8.32]
    step_fill = [SOFT_TEAL, SOFT_ORANGE, SOFT_BLUE]
    step_color = [TEAL, ORANGE, BLUE]
    step_titles = [
        ("A", "多模型基线复现", "FSD、Stay-Positive、LVLM 放到同一实验框架下对比。"),
        ("B", "ADM困难场景分析", "围绕阈值漂移、分数分布与样本级冲突定位失败模式。"),
        ("C", "两阶段联合验证", "把 Stay-Positive 离线分数和 LVLM 语义标签接入训练计算图。"),
    ]
    for i, (x, fillc, sc) in enumerate(zip(x_positions, step_fill, step_color)):
        add_panel(slide, x, 1.82, 3.0, 2.35, fill=fillc, line=fillc, radius=True)
        add_textbox(slide, x + 0.18, 2.02, 0.55, 0.45, step_titles[i][0], font_size=24, color=sc, bold=True, font_face=FONT_LATIN)
        add_textbox(slide, x + 0.18, 2.42, 2.45, 0.4, step_titles[i][1], font_size=18, color=TEXT, bold=True, font_face=FONT_CJK_BOLD)
        add_multiline_text(slide, x + 0.18, 2.8, 2.45, 1.0, [step_titles[i][2]], font_size=14, color=MUTED, font_face=FONT_CJK)
    # arrows
    for x in [3.92, 7.69]:
        arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.58), Inches(0.42), Inches(0.34))
        arr.fill.solid()
        arr.fill.fore_color.rgb = PALE
        arr.line.fill.background()
    # innovation list
    add_panel(slide, 0.78, 4.42, 11.8, 1.95, fill=LIGHT_BG, line=LINE, radius=True)
    add_textbox(slide, 1.0, 4.62, 2.2, 0.3, "创新点 / 贡献", font_size=14, color=TEAL, bold=True, font_face=FONT_CJK_BOLD)
    bullet_lines = [
        "1. 将三类方法纳入统一的多模型分析链条，而不是只做单一模型复现。",
        "2. 以 ADM 为中心做阈值漂移、分数分布、样本级冲突的细粒度分析。",
        "3. 验证 LVLM 结构化语义标签可以作为轻量辅助监督进入训练流程。",
    ]
    add_multiline_text(slide, 1.0, 5.0, 10.7, 1.1, bullet_lines, font_size=18, color=TEXT, font_face=FONT_CJK)

    # Slide 4 baseline chart
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_content_header(
        slide,
        "03 多模型基线复现",
        "多模型基线复现结果",
        "FSD 在多生成器上表现稳定，但不同生成器之间仍存在明显差异。",
        page=4,
    )
    chart_data = CategoryChartData()
    chart_data.categories = ["Midjourney", "SD", "ADM", "BigGAN", "GLIDE", "VQDM"]
    chart_data.add_series("Accuracy", (79.56, 88.34, 75.41, 79.27, 96.67, 75.47))
    chart_data.add_series("AP", (82.04, 91.30, 79.34, 82.40, 96.82, 77.15))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.72),
        Inches(1.55),
        Inches(7.1),
        Inches(4.95),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 20
    chart.category_axis.reverse_order = True
    chart.chart_title.has_text_frame = False
    chart.font.size = Pt(11)
    # series colors
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = TEAL
    # side insights
    add_panel(slide, 8.15, 1.62, 4.45, 4.82, fill=LIGHT_BG, line=LINE, radius=True)
    add_textbox(slide, 8.42, 1.92, 1.8, 0.3, "关键观察", font_size=15, color=TEAL, bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 8.42, 2.35, 3.6, 0.8, "1. GLIDE 最容易判别，Accuracy / AP 接近 97%。", font_size=16, color=TEXT, bold=True, font_face=FONT_CJK)
    add_textbox(slide, 8.42, 3.08, 3.6, 0.8, "2. ADM、VQDM 的指标明显更低，是典型的泛化挑战场景。", font_size=16, color=TEXT, bold=True, font_face=FONT_CJK)
    add_textbox(slide, 8.42, 3.86, 3.6, 0.8, "3. Stay-Positive 在固定真假对照任务中也能保持很高精度，更偏向稳健判别。", font_size=16, color=TEXT, bold=True, font_face=FONT_CJK)
    add_panel(slide, 8.42, 4.86, 3.72, 1.06, fill=SOFT_ORANGE, line=SOFT_ORANGE, radius=True)
    add_textbox(slide, 8.63, 5.03, 3.3, 0.64, "结论：FSD 适合作为未知生成器场景的研究主干，但还需要进一步做困难样本分析和联合验证。", font_size=15, color=TEXT, bold=True, font_face=FONT_CJK)

    # Slide 5 difficult cases
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT_BG)
    add_content_header(
        slide,
        "04 ADM困难场景分析",
        "ADM 场景中的典型失败模式",
        "从高优先级样本池中挑选 3 个典型案例，观察模型为什么会在边界样本上出现分歧。",
        page=5,
    )
    img_dir = ROOT / "analysis" / "adm_priority3_images"
    img_paths = [
        img_dir / "adm_case_001_284_adm_91.PNG",
        img_dir / "adm_case_002_280_adm_7.PNG",
        img_dir / "adm_case_004_385_adm_153.PNG",
    ]
    captions = [
        ("案例 1", "边界样本：人像与前景主体的局部边缘较容易触发判别波动。"),
        ("案例 2", "冲突样本：多主体叠加后，轮廓与纹理会出现明显混叠。"),
        ("案例 3", "稳定外观：整体较自然，但局部细节仍可能诱发模型分歧。"),
    ]
    for i, (p, cap) in enumerate(zip(img_paths, captions)):
        left = 0.8 + i * 4.1
        add_panel(slide, left, 1.78, 3.5, 3.5, fill=WHITE, line=LINE, radius=True)
        add_image(slide, p, left + 0.1, 1.88, 3.3, 2.6)
        add_textbox(slide, left + 0.15, 4.58, 0.9, 0.25, cap[0], font_size=14, color=TEAL if i == 0 else ORANGE if i == 1 else BLUE, bold=True, font_face=FONT_CJK_BOLD)
        add_multiline_text(slide, left + 0.15, 4.9, 3.05, 0.6, [cap[1]], font_size=13, color=TEXT, font_face=FONT_CJK)
    add_panel(slide, 0.8, 5.72, 11.75, 0.74, fill=WHITE, line=LINE, radius=True)
    add_textbox(slide, 1.05, 5.95, 11.2, 0.23, "三类失败模式：阈值漂移   |   分数分布重叠   |   样本级冲突", font_size=17, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 1.05, 6.19, 11.2, 0.23, "结论：模型并非“整体不行”，而是会在局部异常和边界样本上出现不稳定翻转。", font_size=14, color=MUTED, align=PP_ALIGN.CENTER, font_face=FONT_CJK)

    # Slide 6 joint verification
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_content_header(
        slide,
        "05 两阶段联合验证",
        "两阶段联合验证：从离线分数到语义监督",
        "第一阶段验证训练链路，第二阶段验证 LVLM 监督能否真正进入训练计算图。",
        page=6,
    )
    # left timeline
    add_panel(slide, 0.75, 1.7, 4.5, 4.95, fill=LIGHT_BG, line=LINE, radius=True)
    add_rect(slide, 1.05, 2.0, 0.18, 4.1, TEAL)
    add_rect(slide, 1.05, 3.92, 0.18, 2.0, BLUE)
    # stage 1
    add_textbox(slide, 1.28, 2.02, 1.5, 0.3, "Stage 1", font_size=16, color=TEAL, bold=True, font_face=FONT_LATIN)
    add_textbox(slide, 1.28, 2.3, 3.6, 0.4, "FSD + Stay-Positive", font_size=20, color=TEXT, bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 1.28, 2.72, 3.7, 1.15, "10000 step\nADM / SD / Midjourney：95.50 / 97.50\n95.34 / 97.90，87.00 / 90.47\n结论：链路跑通，但有效 SP 样本仍然偏少。", font_size=13, color=TEXT, font_face=FONT_CJK)
    # stage 2
    add_textbox(slide, 1.28, 4.02, 1.5, 0.3, "Stage 2", font_size=16, color=BLUE, bold=True, font_face=FONT_LATIN)
    add_textbox(slide, 1.28, 4.3, 3.6, 0.4, "FSD + LVLM", font_size=20, color=TEXT, bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 1.28, 4.72, 3.7, 1.4, "10000 step\nADM / SD / Midjourney：95.22 / 97.45\n95.33 / 97.72，86.97 / 90.38\nADM 上额外输出 LVLM F1 = 0.1778\n结论：语义监督已进入训练图，但仍需更多标签和更稳的权重策略。", font_size=12, color=TEXT, font_face=FONT_CJK)
    # right evidence table
    add_panel(slide, 5.55, 1.7, 7.05, 4.95, fill=WHITE, line=LINE, radius=True)
    add_textbox(slide, 5.8, 1.93, 2.1, 0.25, "关键证据", font_size=15, color=TEAL, bold=True, font_face=FONT_CJK_BOLD)
    add_panel(slide, 5.8, 2.25, 6.55, 0.78, fill=SOFT_TEAL, line=SOFT_TEAL, radius=True)
    add_textbox(slide, 6.0, 2.48, 6.1, 0.25, "valid_lvlm_samples > 0 说明 LVLM 不只是解释材料，而是真的进入了训练计算图。", font_size=15, color=TEXT, bold=True, font_face=FONT_CJK)
    # metric strip
    metrics = [("6094", "steps_with_valid_lvlm"), ("1.6343", "avg_valid_lvlm_samples_per_step"), ("0.1778", "LVLM F1")]
    for i, (num, label) in enumerate(metrics):
        x = 5.9 + i * 2.12
        add_panel(slide, x, 3.35, 1.9, 1.15, fill=LIGHT_BG, line=LINE, radius=True)
        add_textbox(slide, x + 0.1, 3.56, 1.7, 0.3, num, font_size=24, color=BLUE if i == 2 else TEAL, bold=True, align=PP_ALIGN.CENTER, font_face=FONT_LATIN)
        add_textbox(slide, x + 0.1, 3.95, 1.7, 0.28, label, font_size=11, color=MUTED, align=PP_ALIGN.CENTER, font_face=FONT_LATIN)
    # comparison table
    rows = [
        ["训练设置", "ADM", "SD", "Midjourney", "辅助指标"],
        ["第一阶段\nFSD + Stay-Positive", "95.50 / 97.50", "95.34 / 97.90", "87.00 / 90.47", "SP 有效样本较少"],
        ["第二阶段\nFSD + LVLM", "95.22 / 97.45", "95.33 / 97.72", "86.97 / 90.38", "LVLM F1 = 0.1778"],
    ]
    table = slide.shapes.add_table(3, 5, Inches(5.84), Inches(4.72), Inches(6.6), Inches(1.45)).table
    col_widths = [1.25, 1.5, 1.3, 1.45, 1.1]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    for r in range(3):
        for c in range(5):
            cell = table.cell(r, c)
            cell.text = rows[r][c]
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.name = FONT_CJK
                    run.font.size = Pt(12 if r else 11)
                    run.font.bold = True if (r == 0 or c == 4) else False
                    run.font.color.rgb = WHITE if r == 0 else TEXT
    add_textbox(slide, 5.85, 6.34, 6.6, 0.22, "结论：第二阶段属于“最小量化验证”，它证明了语义监督已经接入，但还没有显著抬升主任务指标。", font_size=13, color=MUTED, align=PP_ALIGN.CENTER, font_face=FONT_CJK)

    # Slide 7 summary and outlook
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT_BG)
    add_content_header(
        slide,
        "06 总结与展望",
        "研究结论与后续工作",
        "本文的价值在于把三个层次的证据整合到同一个研究过程中：指标、失败模式和联合验证。",
        page=7,
    )
    add_panel(slide, 0.75, 1.72, 5.95, 4.95, fill=WHITE, line=LINE, radius=True)
    add_panel(slide, 6.9, 1.72, 5.68, 4.95, fill=WHITE, line=LINE, radius=True)
    add_textbox(slide, 1.0, 1.98, 2.0, 0.25, "主要结论", font_size=15, color=TEAL, bold=True, font_face=FONT_CJK_BOLD)
    concl = [
        "1. 完成了 FSD、Stay-Positive、LVLM 的统一复现与对比。",
        "2. 证明 ADM 的主要难点来自阈值漂移和样本级冲突，而不只是整体精度不足。",
        "3. 验证了 LVLM 语义标签可以作为轻量辅助监督接入训练流程。",
    ]
    add_multiline_text(slide, 1.0, 2.32, 5.2, 1.6, concl, font_size=17, color=TEXT, font_face=FONT_CJK)
    add_rule(slide, 1.0, 4.3, 5.0, PALE, weight=2)
    add_textbox(slide, 1.0, 4.48, 2.0, 0.25, "后续工作", font_size=15, color=ORANGE, bold=True, font_face=FONT_CJK_BOLD)
    future = [
        "• 扩大 LVLM 标签规模和类别均衡性。",
        "• 优化阈值校准、样本筛选和损失权重。",
        "• 扩展到更多生成器和更多随机种子消融。",
    ]
    add_multiline_text(slide, 1.0, 4.78, 5.1, 1.45, future, font_size=16, color=TEXT, font_face=FONT_CJK)
    add_textbox(slide, 7.2, 1.98, 2.0, 0.25, "答辩时可这样总结", font_size=15, color=BLUE, bold=True, font_face=FONT_CJK_BOLD)
    summary_box = [
        "这篇论文不是单纯复现已有方法，",
        "而是把多模型基线比较、ADM 困难场景剖析、",
        "两阶段联合验证串成了一条完整链条。",
        "",
        "它说明：在 AI 伪造图像识别中，",
        "“更稳的判别”与“更可解释的监督”",
        "可以被放到同一个研究框架里讨论。",
    ]
    add_multiline_text(slide, 7.2, 2.32, 4.95, 2.65, summary_box, font_size=18, color=TEXT, bold=False, font_face=FONT_CJK)
    add_panel(slide, 7.15, 5.22, 4.95, 0.9, fill=SOFT_BLUE, line=SOFT_BLUE, radius=True)
    add_textbox(slide, 7.38, 5.47, 4.5, 0.25, "谢谢各位老师，请批评指正。", font_size=19, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font_face=FONT_CJK_BOLD)
    add_footer(slide, 7)

    # Slide 8 Q&A
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, NAVY)
    add_rect(slide, 0.0, 0.0, 13.333, 7.5, NAVY)
    add_rect(slide, 0.0, 0.0, 0.22, 7.5, ORANGE)
    add_rect(slide, 9.85, 1.2, 2.7, 0.18, TEAL)
    add_rect(slide, 9.85, 1.6, 2.1, 0.18, ORANGE)
    add_textbox(slide, 0.82, 1.55, 5.8, 0.55, "答疑交流", font_size=18, color=RGBColor(184, 205, 230), bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 0.82, 2.15, 6.5, 0.95, "谢谢各位老师", font_size=36, color=WHITE, bold=True, font_face=FONT_CJK_BOLD)
    add_textbox(slide, 0.82, 3.15, 6.8, 0.6, "恳请批评指正", font_size=22, color=RGBColor(214, 231, 242), font_face=FONT_CJK)
    add_rule(slide, 0.82, 3.98, 1.7, TEAL, weight=4)
    add_textbox(slide, 0.82, 4.28, 7.6, 0.95, "如需进一步展开，可围绕：\n1. 多模型基线差异  2. ADM 困难样本  3. LVLM 联合验证结果", font_size=16, color=RGBColor(215, 227, 239), font_face=FONT_CJK)
    add_textbox(slide, 9.95, 3.55, 2.05, 0.72, "Q&A", font_size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font_face=FONT_LATIN)

    return prs


def draw_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, font: ImageFont.FreeTypeFont, fill, max_width: int | None = None, line_gap: int = 8):
    x, y = xy
    if max_width is None:
        draw.multiline_text((x, y), text, font=font, fill=fill, spacing=line_gap)
        bbox = draw.multiline_textbbox((x, y), text, font=font, spacing=line_gap)
        return bbox
    lines: list[str] = []
    for paragraph in text.split("\n"):
        if paragraph == "":
            lines.append("")
            continue
        current = ""
        for ch in paragraph:
            test = current + ch
            if draw.textbbox((0, 0), test, font=font)[2] <= max_width or not current:
                current = test
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    draw.multiline_text((x, y), "\n".join(lines), font=font, fill=fill, spacing=line_gap)
    return draw.multiline_textbbox((x, y), "\n".join(lines), font=font, spacing=line_gap)


def draw_round_rect(draw: ImageDraw.ImageDraw, box, fill, outline=None, width=2, radius=22):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def render_preview_slide_1(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (15, 35, 61))
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, 28, SLIDE_H), fill=(244, 160, 72))
    d.rectangle((980, 60, 1235, 650), fill=(19, 47, 82))
    d.rectangle((1010, 100, 1200, 115), fill=(41, 173, 163))
    d.rectangle((1010, 145, 1160, 160), fill=(244, 160, 72))
    d.rectangle((1010, 190, 1120, 205), fill=(52, 121, 227))
    f1 = load_font(52, bold=True)
    f2 = load_font(28, bold=False)
    f3 = load_font(20, bold=False)
    f4 = load_font(15, bold=True)
    d.text((62, 56), "答辩汇报", font=f4, fill=(180, 205, 230))
    d.text((62, 98), "AI伪造图像识别研究", font=f1, fill=(255, 255, 255))
    d.text((62, 190), "面向泛化检测的多模型识别与联合验证", font=f2, fill=(214, 231, 242))
    d.rectangle((62, 270, 260, 276), fill=(41, 173, 163))
    d.text((62, 294), "从基线复现，到困难场景分析，再到两阶段联合验证", font=f3, fill=(207, 219, 233))
    # chips
    chips = [("FSD", (41, 173, 163), 62), ("Stay-Positive", (244, 160, 72), 155), ("LVLM", (52, 121, 227), 335)]
    for txt, fill, x in chips:
        draw_round_rect(d, (x, 378, x + (80 if txt=="FSD" else 140 if txt=="LVLM" else 160), 412), fill=fill, radius=14)
        d.text((x + 18, 385), txt, font=load_font(16, bold=True), fill=(255,255,255))
    d.text((62, 470), "学生：待填写\n指导教师：待填写\n完成时间：2026年5月", font=f3, fill=(215, 227, 239), spacing=8)
    d.text((1015, 348), "多模型\n基线复现", font=load_font(30, bold=True), fill=(255,255,255), align="center")
    d.text((1005, 455), "困难场景\n样本分析", font=load_font(30, bold=True), fill=(255,255,255))
    d.text((1010, 560), "联合验证\n与总结", font=load_font(30, bold=True), fill=(255,255,255))
    img.save(path)


def render_preview_slide_2(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (248, 250, 252))
    d = ImageDraw.Draw(img)
    d.text((62, 24), "01 背景与问题", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((62, 56), "为什么 AI 伪造图像检测越来越难", font=load_font(34, bold=True), fill=(26,34,48))
    d.rectangle((62, 108, 220, 112), fill=(41,173,163))
    d.text((62, 122), "生成器迭代快，单一检测器很难同时兼顾泛化性、稳定性与解释性。", font=load_font(18), fill=(95,107,123))
    panel_specs = [
        (62, 170, (41, 173, 163), "01", "跨生成器泛化不足", [
            "训练时见过的生成器，",
            "和部署时遇到的生成器往往不一致。",
            "单一模型容易学到局部模式。",
        ]),
        (435, 170, (244, 160, 72), "02", "阈值与分数分布不稳定", [
            "真伪样本在边界区域容易重叠，",
            "0.5 附近最容易翻转。",
            "不同协议下阈值会漂移。",
        ]),
        (808, 170, (52, 121, 227), "03", "复杂样本难解释", [
            "误判样本往往不是“完全不像”，",
            "而是局部冲突和混合证据。",
            "仅给真假结果不够。",
        ]),
    ]
    for x, y, c, num, h, lines in panel_specs:
        draw_round_rect(d, (x, y, x + 330, y + 420), fill=(255,255,255), outline=(214,221,231), radius=20)
        d.text((x+18, y+22), num, font=load_font(30, bold=True), fill=c)
        d.text((x+18, y+68), h, font=load_font(22, bold=True), fill=(26,34,48))
        d.rectangle((x+18, y+110, x+90, y+114), fill=c)
        draw_text(d, (x+18, y+140), "\n".join(lines), font=load_font(17), fill=(26,34,48), max_width=280, line_gap=6)
    draw_round_rect(d, (62, 595, 1213, 641), fill=(235,250,248), outline=(235,250,248), radius=16)
    d.text((92, 609), "研究目标：构建一个面向泛化检测的多模型协同分析框架，把“复现、分析、联合验证”串成一条可复核的研究链条。", font=load_font(17, bold=True), fill=(26,34,48))
    img.save(path)


def render_preview_slide_3(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (255,255,255))
    d = ImageDraw.Draw(img)
    d.text((62, 24), "02 研究思路与创新点", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((62, 56), "本文的研究主线与创新点", font=load_font(34, bold=True), fill=(26,34,48))
    d.rectangle((62, 108, 220, 112), fill=(41,173,163))
    d.text((62, 122), "不是只做一个结果，而是把复现、困难场景分析和联合验证连成闭环。", font=load_font(18), fill=(95,107,123))
    steps = [
        (62, 170, (235,250,248), (41,173,163), "A", "多模型基线复现", "FSD、Stay-Positive、LVLM 放到同一实验框架下对比。"),
        (441, 170, (255,245,232), (244,160,72), "B", "ADM困难场景分析", "围绕阈值漂移、分数分布与样本级冲突定位失败模式。"),
        (820, 170, (237,244,255), (52,121,227), "C", "两阶段联合验证", "把离线分数和 LVLM 语义标签接入训练计算图。"),
    ]
    for x, y, fill, c, letter, title, body in steps:
        draw_round_rect(d, (x, y, x+315, y+220), fill=fill, outline=fill, radius=20)
        d.text((x+18, y+22), letter, font=load_font(28, bold=True), fill=c)
        d.text((x+18, y+70), title, font=load_font(21, bold=True), fill=(26,34,48))
        d.text((x+18, y+112), body, font=load_font(16), fill=(95,107,123))
    # arrows
    d.polygon([(388, 272), (422, 272), (422, 263), (448, 283), (422, 303), (422, 294), (388, 294)], fill=(233,239,245))
    d.polygon([(767, 272), (801, 272), (801, 263), (827, 283), (801, 303), (801, 294), (767, 294)], fill=(233,239,245))
    draw_round_rect(d, (62, 390, 1210, 583), fill=(248,250,252), outline=(214,221,231), radius=20)
    d.text((92, 414), "创新点 / 贡献", font=load_font(16, bold=True), fill=(41,173,163))
    lines = [
        "1. 将三类方法纳入统一的多模型分析链条，而不是只做单一模型复现。",
        "2. 以 ADM 为中心做阈值漂移、分数分布、样本级冲突的细粒度分析。",
        "3. 验证 LVLM 结构化语义标签可以作为轻量辅助监督进入训练流程。",
    ]
    draw_text(d, (92, 455), "\n".join(lines), font=load_font(22), fill=(26,34,48), max_width=1080, line_gap=12)
    img.save(path)


def render_preview_slide_4(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (255,255,255))
    d = ImageDraw.Draw(img)
    d.text((62, 24), "03 多模型基线复现", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((62, 56), "多模型基线复现结果", font=load_font(34, bold=True), fill=(26,34,48))
    d.rectangle((62, 108, 220, 112), fill=(41,173,163))
    d.text((62, 122), "FSD 在多生成器上表现稳定，但不同生成器之间仍存在明显差异。", font=load_font(18), fill=(95,107,123))
    # chart area
    draw_round_rect(d, (62, 170, 770, 640), fill=(255,255,255), outline=(214,221,231), radius=20)
    # simple clustered bars
    cats = ["Midjourney", "SD", "ADM", "BigGAN", "GLIDE", "VQDM"]
    acc = [79.56, 88.34, 75.41, 79.27, 96.67, 75.47]
    ap = [82.04, 91.30, 79.34, 82.40, 96.82, 77.15]
    y0 = 215
    for idx, cat in enumerate(cats):
        y = y0 + idx * 64
        d.text((78, y), cat, font=load_font(16, bold=False), fill=(26,34,48))
        # bar background
        d.rectangle((200, y+2, 720, y+18), fill=(237,242,248))
        d.rectangle((200, y+2, 200 + int(5.1*acc[idx]), y+18), fill=(52,121,227))
        d.rectangle((200, y+24, 720, y+40), fill=(237,242,248))
        d.rectangle((200, y+24, 200 + int(5.1*ap[idx]), y+40), fill=(41,173,163))
        d.text((730, y-2), f"{acc[idx]:.2f}", font=load_font(14), fill=(26,34,48))
        d.text((730, y+22), f"{ap[idx]:.2f}", font=load_font(14), fill=(26,34,48))
    d.text((206, 190), "Accuracy", font=load_font(14, bold=True), fill=(52,121,227))
    d.text((320, 190), "AP", font=load_font(14, bold=True), fill=(41,173,163))
    # right panel
    draw_round_rect(d, (807, 170, 1210, 640), fill=(248,250,252), outline=(214,221,231), radius=20)
    d.text((836, 198), "关键观察", font=load_font(16, bold=True), fill=(41,173,163))
    obs = [
        "1. GLIDE 最容易判别，Accuracy / AP 接近 97%。",
        "2. ADM、VQDM 是典型的泛化挑战场景。",
        "3. Stay-Positive 在固定真假对照任务中也能保持很高精度。",
    ]
    draw_text(d, (836, 250), "\n\n".join(obs), font=load_font(18), fill=(26,34,48), max_width=330, line_gap=8)
    draw_round_rect(d, (836, 487, 1180, 585), fill=(255,245,232), outline=(255,245,232), radius=18)
    draw_text(d, (852, 505), "结论：FSD 适合作为未知生成器场景的研究主干，但还需要进一步做困难样本分析和联合验证。", font=load_font(17, bold=True), fill=(26,34,48), max_width=310, line_gap=6)
    img.save(path)


def render_preview_slide_5(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (248,250,252))
    d = ImageDraw.Draw(img)
    d.text((62, 24), "04 ADM困难场景分析", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((62, 56), "ADM 场景中的典型失败模式", font=load_font(34, bold=True), fill=(26,34,48))
    d.rectangle((62, 108, 220, 112), fill=(41,173,163))
    d.text((62, 122), "从高优先级样本池中挑选 3 个典型案例，观察模型为什么会在边界样本上出现分歧。", font=load_font(18), fill=(95,107,123))
    img_dir = ROOT / "analysis" / "adm_priority3_images"
    paths = [
        img_dir / "adm_case_001_284_adm_91.PNG",
        img_dir / "adm_case_002_280_adm_7.PNG",
        img_dir / "adm_case_004_385_adm_153.PNG",
    ]
    infos = [
        ("案例 1", "边界样本：人像与前景主体的局部边缘较容易触发判别波动。"),
        ("案例 2", "冲突样本：多主体叠加后，轮廓与纹理会出现明显混叠。"),
        ("案例 3", "稳定外观：整体较自然，但局部细节仍可能诱发模型分歧。"),
    ]
    colors = [(41,173,163), (244,160,72), (52,121,227)]
    for i, (p, info, c) in enumerate(zip(paths, infos, colors)):
        x = 64 + i*377
        draw_round_rect(d, (x, 170, x+342, 520), fill=(255,255,255), outline=(214,221,231), radius=18)
        img_case = Image.open(p).convert("RGB").resize((322, 250))
        img.paste(img_case, (x+10, 180))
        d.text((x+16, 454), info[0], font=load_font(15, bold=True), fill=c)
        draw_text(d, (x+16, 482), info[1], font=load_font(13), fill=(26,34,48), max_width=310, line_gap=4)
    draw_round_rect(d, (64, 554, 1194, 624), fill=(255,255,255), outline=(214,221,231), radius=16)
    d.text((82, 577), "三类失败模式：阈值漂移   |   分数分布重叠   |   样本级冲突", font=load_font(20, bold=True), fill=(26,34,48))
    d.text((82, 603), "结论：模型并非“整体不行”，而是会在局部异常和边界样本上出现不稳定翻转。", font=load_font(15), fill=(95,107,123))
    img.save(path)


def render_preview_slide_6(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (255,255,255))
    d = ImageDraw.Draw(img)
    d.text((62, 24), "05 两阶段联合验证", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((62, 56), "两阶段联合验证：从离线分数到语义监督", font=load_font(34, bold=True), fill=(26,34,48))
    d.rectangle((62, 108, 265, 112), fill=(41,173,163))
    d.text((62, 122), "第一阶段验证训练链路，第二阶段验证 LVLM 监督能否真正进入训练计算图。", font=load_font(18), fill=(95,107,123))
    draw_round_rect(d, (62, 170, 488, 640), fill=(248,250,252), outline=(214,221,231), radius=18)
    draw_round_rect(d, (509, 170, 1210, 640), fill=(255,255,255), outline=(214,221,231), radius=18)
    d.rectangle((85, 200, 95, 565), fill=(41,173,163))
    d.rectangle((85, 420, 95, 565), fill=(52,121,227))
    d.text((110, 205), "Stage 1", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((110, 235), "FSD + Stay-Positive", font=load_font(22, bold=True), fill=(26,34,48))
    draw_text(d, (110, 275), "10000 step\nADM / SD / Midjourney：95.50 / 97.50\n95.34 / 97.90，87.00 / 90.47\n结论：链路跑通，但有效 SP 样本仍然偏少。", font=load_font(14), fill=(26,34,48), max_width=335, line_gap=6)
    d.text((110, 430), "Stage 2", font=load_font(16, bold=True), fill=(52,121,227))
    d.text((110, 460), "FSD + LVLM", font=load_font(22, bold=True), fill=(26,34,48))
    draw_text(d, (110, 500), "10000 step\nADM / SD / Midjourney：95.22 / 97.45\n95.33 / 97.72，86.97 / 90.38\nADM 上额外输出 LVLM F1 = 0.1778\n结论：语义监督已进入训练图，但仍需更多标签和更稳的权重策略。", font=load_font(13), fill=(26,34,48), max_width=335, line_gap=5)
    d.text((535, 200), "关键证据", font=load_font(16, bold=True), fill=(41,173,163))
    draw_round_rect(d, (535, 236, 1178, 300), fill=(235,250,248), outline=(235,250,248), radius=16)
    d.text((552, 256), "valid_lvlm_samples > 0 说明 LVLM 不只是解释材料，而是真的进入了训练计算图。", font=load_font(15, bold=True), fill=(26,34,48))
    # metrics strip
    metrics = [("6094", "steps_with_valid_lvlm"), ("1.6343", "avg_valid_lvlm_samples_per_step"), ("0.1778", "LVLM F1")]
    for i, (num, label) in enumerate(metrics):
        x = 535 + i * 215
        draw_round_rect(d, (x, 330, x+195, 420), fill=(248,250,252), outline=(214,221,231), radius=16)
        d.text((x+98, 352), num, font=load_font(24, bold=True), fill=(52,121,227) if i == 2 else (41,173,163), anchor="ma")
        d.text((x+98, 388), label, font=load_font(12), fill=(95,107,123), anchor="ma")
    # table
    headers = ["训练设置", "ADM", "SD", "Midjourney", "辅助指标"]
    rows = [
        ["第一阶段\nFSD + Stay-Positive", "95.50 / 97.50", "95.34 / 97.90", "87.00 / 90.47", "SP 有效样本较少"],
        ["第二阶段\nFSD + LVLM", "95.22 / 97.45", "95.33 / 97.72", "86.97 / 90.38", "LVLM F1 = 0.1778"],
    ]
    x0, y0 = 535, 445
    widths = [158, 142, 132, 145, 128]
    h = 56
    x = x0
    for i, head in enumerate(headers):
        draw_round_rect(d, (x, y0, x+widths[i], y0+h), fill=(15,35,61), outline=(15,35,61), radius=8)
        d.text((x+widths[i]/2, y0+18), head, font=load_font(12, bold=True), fill=(255,255,255), anchor="ma")
        x += widths[i]
    for r, row in enumerate(rows):
        y = y0 + h + r*h
        x = x0
        fill = (255,255,255) if r == 0 else (248,250,252)
        for i, cell in enumerate(row):
            draw_round_rect(d, (x, y, x+widths[i], y+h), fill=fill, outline=(214,221,231), radius=8)
            d.text((x+widths[i]/2, y+18), cell, font=load_font(11, bold=(i==4)), fill=(26,34,48), anchor="ma")
            x += widths[i]
    d.text((536, 618), "结论：第二阶段属于“最小量化验证”，它证明了语义监督已经接入，但还没有显著抬升主任务指标。", font=load_font(13), fill=(95,107,123))
    img.save(path)


def render_preview_slide_7(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (248,250,252))
    d = ImageDraw.Draw(img)
    d.text((62, 24), "06 总结与展望", font=load_font(16, bold=True), fill=(41,173,163))
    d.text((62, 56), "研究结论与后续工作", font=load_font(34, bold=True), fill=(26,34,48))
    d.rectangle((62, 108, 220, 112), fill=(41,173,163))
    d.text((62, 122), "本文的价值在于把三个层次的证据整合到同一个研究过程中：指标、失败模式和联合验证。", font=load_font(18), fill=(95,107,123))
    draw_round_rect(d, (62, 170, 620, 640), fill=(255,255,255), outline=(214,221,231), radius=18)
    draw_round_rect(d, (650, 170, 1210, 640), fill=(255,255,255), outline=(214,221,231), radius=18)
    d.text((90, 198), "主要结论", font=load_font(16, bold=True), fill=(41,173,163))
    conclusion = [
        "1. 完成了 FSD、Stay-Positive、LVLM 的统一复现与对比。",
        "2. 证明 ADM 的主要难点来自阈值漂移和样本级冲突。",
        "3. 验证了 LVLM 语义标签可以作为轻量辅助监督接入训练流程。",
    ]
    draw_text(d, (90, 240), "\n\n".join(conclusion), font=load_font(18), fill=(26,34,48), max_width=470, line_gap=8)
    d.rectangle((90, 417, 500, 419), fill=(233,239,245))
    d.text((90, 444), "后续工作", font=load_font(16, bold=True), fill=(244,160,72))
    future = [
        "• 扩大 LVLM 标签规模和类别均衡性。",
        "• 优化阈值校准、样本筛选和损失权重。",
        "• 扩展到更多生成器和更多随机种子消融。",
    ]
    draw_text(d, (90, 484), "\n".join(future), font=load_font(17), fill=(26,34,48), max_width=470, line_gap=6)
    d.text((680, 198), "答辩时可这样总结", font=load_font(16, bold=True), fill=(52,121,227))
    summary = [
        "这篇论文不是单纯复现已有方法，",
        "而是把多模型基线比较、ADM 困难场景剖析、",
        "两阶段联合验证串成了一条完整链条。",
        "",
        "它说明：在 AI 伪造图像识别中，",
        "“更稳的判别”与“更可解释的监督”",
        "可以被放到同一个研究框架里讨论。",
    ]
    draw_text(d, (680, 240), "\n".join(summary), font=load_font(19), fill=(26,34,48), max_width=450, line_gap=8)
    draw_round_rect(d, (680, 520, 1140, 585), fill=(237,244,255), outline=(237,244,255), radius=16)
    d.text((910, 551), "谢谢各位老师，请批评指正。", font=load_font(20, bold=True), fill=(26,34,48), anchor="ma")
    img.save(path)


def render_preview_slide_8(path: Path):
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), (15, 35, 61))
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, 28, SLIDE_H), fill=(244, 160, 72))
    d.rectangle((945, 115, 1195, 131), fill=(41, 173, 163))
    d.rectangle((945, 154, 1145, 170), fill=(244, 160, 72))
    d.text((70, 136), "答疑交流", font=load_font(20, bold=True), fill=(184, 205, 230))
    d.text((70, 192), "谢谢各位老师", font=load_font(54, bold=True), fill=(255, 255, 255))
    d.text((70, 292), "恳请批评指正", font=load_font(30), fill=(214, 231, 242))
    d.rectangle((70, 383, 265, 389), fill=(41, 173, 163))
    d.text((70, 420), "如需进一步展开，可围绕：\n1. 多模型基线差异  2. ADM 困难样本\n3. LVLM 联合验证结果", font=load_font(21), fill=(215, 227, 239), spacing=8)
    d.text((980, 338), "Q&A", font=load_font(56, bold=True), fill=(255, 255, 255))
    img.save(path)


def render_previews():
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    render_preview_slide_1(PREVIEW_DIR / "slide-01.png")
    render_preview_slide_2(PREVIEW_DIR / "slide-02.png")
    render_preview_slide_3(PREVIEW_DIR / "slide-03.png")
    render_preview_slide_4(PREVIEW_DIR / "slide-04.png")
    render_preview_slide_5(PREVIEW_DIR / "slide-05.png")
    render_preview_slide_6(PREVIEW_DIR / "slide-06.png")
    render_preview_slide_7(PREVIEW_DIR / "slide-07.png")
    render_preview_slide_8(PREVIEW_DIR / "slide-08.png")


def main():
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = build_pptx()
    prs.save(OUT_PPTX)
    render_previews()
    print(f"pptx_saved={OUT_PPTX}")
    print(f"preview_dir={PREVIEW_DIR}")


if __name__ == "__main__":
    main()
